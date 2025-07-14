import pptx
import glob
from pathlib import Path
import sys
from PyQt5.QtWidgets import QApplication, QFileDialog


def choose_folder_qt():
    '''
    Opens system dialogues to find the input and output directories
    '''
    # Initiates the PyQt object needed to use system diaglogues
    app = QApplication(sys.argv)

    # Ask for and store the input folder
    in_folder = QFileDialog.getExistingDirectory(
        None, "Select folder containing Powerpoint pptx files")

    # Ask for and store the output folder
    out_folder = QFileDialog.getExistingDirectory(
        None, "Select or create a different folder for the clean files")

    # Check both have been specified
    assert in_folder and out_folder, "Input and output folders \
        must be specified"
    return in_folder, out_folder


def remove_notes(input_path, output_path):
    '''
    Remove the notes pages, slide by slide
    '''
    # Read the powerpoint presentation
    prs = pptx.Presentation(input_path)
    x = 0
    # Iterate through the slides
    for slide in prs.slides:
        # Check slide for notes
        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
            # The notes are contained in "shapes"
            for shape in notes_slide.shapes:
                # Clear shapes with text frames
                if shape.has_text_frame:
                    shape.text_frame.clear()
                    x += 1
    # Output whether anything was removed
    if x > 0:
        print(f"Removed notes from {input_path}")
    else:
        print(f"{input_path} had no notes")
    # Save the resulting file
    prs.save(output_path)


def main():
    '''
    Main code execution
    '''
    # Run choose_folder_qt to find the input and output folder
    infold, outfold = choose_folder_qt()
    # Iteratre through the files ending in pptx in the input folder
    for infile in glob.glob(f"{infold}/*pptx"):
        # Determine the stem of the filename
        # e.g. xxxxx/hello.txt would be hello
        file_pref = Path(infile).stem
        # Build the name for the output file
        outfile = f"{outfold}/{file_pref}.pptx"
        # Use the remove_notes function to remove the notes and save the
        # resulting file
        remove_notes(infile, outfile)


if __name__ == "__main__":
    main()
